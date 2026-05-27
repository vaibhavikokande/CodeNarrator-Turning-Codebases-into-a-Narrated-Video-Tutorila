"""
export_pptx.py — Educational PowerPoint Export (REDESIGNED).

Creates a proper educational .pptx from the generated markdown chapters.
6 slides per chapter:
  1. Title Slide          — chapter name, repo, date
  2. Overview             — "What is [topic]?" + key bullets + AI speaker notes
  3. Code Walkthrough     — actual code snippet + line explanations + AI notes
  4. Key Concepts         — concept grid with definitions + analogies + AI notes
  5. Tutorial Steps       — numbered how-to + tips + AI notes
  6. Summary & Next       — recap checklist + next chapter link

Design: clean dark-navy + white educational theme (blue accent).
Speaker notes on EVERY slide — generated via the existing LLM router.

Only modifies: this file + the /export/pptx endpoint in main.py.
All other pipeline code is untouched.
"""

import asyncio
import logging
import re
import textwrap
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ── Color palette ─────────────────────────────────────────────────────────────
# Navy-blue educational theme
_C_NAVY     = (30,  58,  95)   # #1E3A5F — dark navy background
_C_BLUE     = (37,  99,  235)  # #2563EB — medium blue accent
_C_LBLUE    = (219, 234, 254)  # #DBEAFE — light blue tint
_C_WHITE    = (255, 255, 255)
_C_OFFWHITE = (248, 250, 252)  # #F8FAFC — slide background
_C_DARK     = (15,  23,  42)   # #0F172A — body text
_C_MUTED    = (100, 116, 139)  # #64748B — muted text
_C_CODE_BG  = (15,  23,  42)   # #0F172A — code block bg
_C_CODE_TXT = (226, 232, 240)  # #E2E8F0 — code text
_C_GREEN    = (16,  185, 129)  # #10B981 — checkmarks
_C_ACCENT   = (99,  102, 241)  # #6366F1 — purple accent (chapter number)

# Slide dimensions — 16:9 widescreen
_W_IN = 13.333
_H_IN = 7.5


# ── Markdown parser ───────────────────────────────────────────────────────────

def _parse_chapter(md_text: str) -> Dict[str, Any]:
    """
    Extract structured data from a chapter markdown file.
    Returns dict with: title, motivation, core_concepts, code_blocks,
                        practical_usage, conclusion, all_text
    """
    lines      = md_text.splitlines()
    title      = ""
    sections: Dict[str, List[str]] = {}
    current    = "intro"
    code_blocks: List[str] = []
    in_code    = False
    code_buf: List[str] = []

    for line in lines:
        # Track code blocks
        if line.strip().startswith("```"):
            if in_code:
                if code_buf:
                    code_blocks.append("\n".join(code_buf))
                code_buf = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buf.append(line)
            continue

        # Chapter title
        if line.startswith("# ") and not title:
            # Strip "Chapter N: " prefix if present
            raw = line[2:].strip()
            raw = re.sub(r"^Chapter\s+\d+:\s*", "", raw, flags=re.IGNORECASE)
            title = raw
            continue

        # Section headings
        if line.startswith("## "):
            heading = line[3:].strip().lower()
            if "motivat" in heading:
                current = "motivation"
            elif "core" in heading or "concept" in heading:
                current = "core_concepts"
            elif "practical" in heading or "usage" in heading:
                current = "practical"
            elif "internal" in heading or "mechanic" in heading:
                current = "internal"
            elif "conclus" in heading:
                current = "conclusion"
            else:
                current = heading.replace(" ", "_")[:20]
            sections.setdefault(current, [])
            continue

        sections.setdefault(current, [])
        sections[current].append(line)

    def _join(key: str) -> str:
        return "\n".join(sections.get(key, [])).strip()

    # Extract bullet points from any section
    def _bullets(text: str, max_n: int = 5) -> List[str]:
        pts = []
        for ln in text.splitlines():
            ln = ln.strip()
            if ln.startswith(("- ", "* ", "• ")):
                pts.append(ln[2:].strip())
            elif re.match(r"^\d+\.\s", ln):
                pts.append(re.sub(r"^\d+\.\s+", "", ln))
        return pts[:max_n]

    motivation   = _join("motivation")
    core         = _join("core_concepts")
    practical    = _join("practical")
    conclusion   = _join("conclusion")

    # Overview bullets: from motivation or first intro lines
    overview_bullets = _bullets(motivation) or _bullets(core)
    if not overview_bullets:
        # Fall back: first 5 non-empty non-heading lines
        for ln in (sections.get("intro", []) + sections.get("motivation", [])):
            ln = ln.strip()
            if ln and not ln.startswith("#") and len(ln) > 20:
                overview_bullets.append(ln[:120])
                if len(overview_bullets) >= 4:
                    break

    # Key concepts: pick bold **Term** patterns
    key_concepts: List[Tuple[str, str]] = []
    for section_text in [core, motivation, practical]:
        for m in re.finditer(r"\*\*(.+?)\*\*\s*[:\-–]?\s*(.+?)(?:\n|$)", section_text):
            name = m.group(1).strip()
            defn = m.group(2).strip()[:100]
            if name and defn and len(name) < 40:
                key_concepts.append((name, defn))
                if len(key_concepts) >= 4:
                    break
        if len(key_concepts) >= 4:
            break

    # Tutorial steps: numbered list from practical usage
    steps = _bullets(practical, max_n=5)
    if not steps:
        steps = _bullets(core, max_n=5)

    # Best code snippet: prefer longer ones
    best_code = ""
    for block in sorted(code_blocks, key=len, reverse=True):
        if len(block.strip()) > 20:
            best_code = block.strip()
            break

    # Conclusion / summary bullets
    summary_bullets = _bullets(conclusion, max_n=4)
    if not summary_bullets:
        # grab last non-empty sentences from conclusion
        for ln in conclusion.splitlines():
            ln = ln.strip()
            if len(ln) > 20:
                summary_bullets.append(ln[:120])
                if len(summary_bullets) >= 3:
                    break

    return {
        "title":            title or "Chapter",
        "motivation":       motivation,
        "core":             core,
        "practical":        practical,
        "conclusion":       conclusion,
        "overview_bullets": overview_bullets,
        "key_concepts":     key_concepts,
        "code_blocks":      code_blocks,
        "best_code":        best_code,
        "steps":            steps,
        "summary_bullets":  summary_bullets,
        "all_text":         md_text[:3000],
    }


# ── LLM speaker notes ─────────────────────────────────────────────────────────

async def _gen_notes(topic: str, context: str, slide_type: str) -> str:
    """Generate 3-5 sentence speaker notes using the existing LLM router."""
    prompts = {
        "overview": (
            f"You are an educational content writer. Explain '{topic}' in 4-5 sentences "
            f"for a complete beginner learning to code. "
            f"Include: what it is, why it matters, and a real-world analogy. "
            f"Context: {context[:500]}"
        ),
        "code": (
            f"You are a coding instructor. Explain this code snippet from '{topic}' "
            f"line by line for a beginner. Include what each important part does and why. "
            f"Keep it to 4-5 sentences. Code context: {context[:400]}"
        ),
        "concepts": (
            f"You are a teacher explaining '{topic}' to a student. "
            f"In 4-5 sentences, explain the key concepts they need to understand. "
            f"Use simple language and concrete examples. Context: {context[:500]}"
        ),
        "steps": (
            f"You are a coding mentor. Give a beginner developer practical advice "
            f"for implementing '{topic}'. Include tips, warnings, and best practices. "
            f"Keep it to 4-5 sentences. Context: {context[:400]}"
        ),
        "summary": (
            f"Summarize what a student should remember after learning about '{topic}'. "
            f"Give 3 key takeaways in 4-5 sentences total. "
            f"End with encouragement to practice. Context: {context[:300]}"
        ),
    }
    prompt = prompts.get(slide_type, prompts["overview"])
    try:
        from llm.router import llm_call
        notes = await llm_call(prompt, bypass_cache=False)
        return notes.strip()[:1200] if notes else ""
    except Exception as exc:
        logger.warning("LLM notes failed for %s/%s: %s", topic, slide_type, exc)
        return f"This slide covers key aspects of {topic}. Review the content carefully and practice implementing the concepts shown."


# ── python-pptx slide builders ────────────────────────────────────────────────

def _rgb(c: Tuple) -> Any:
    from pptx.dml.color import RGBColor
    return RGBColor(*c)


def _inches(val: float) -> Any:
    from pptx.util import Inches
    return Inches(val)


def _pt(val: int) -> Any:
    from pptx.util import Pt
    return Pt(val)


def _add_rect(slide, left: float, top: float, width: float, height: float,
              fill: Tuple, alpha: int = 255) -> Any:
    """Add a filled rectangle to a slide."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    return shape


def _add_textbox(slide, text: str, left: float, top: float,
                 width: float, height: float,
                 font_size: int = 16, bold: bool = False,
                 color: Tuple = _C_DARK, italic: bool = False,
                 word_wrap: bool = True, align_center: bool = False) -> Any:
    """Add a styled text box."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txBox


def _add_bullet_textbox(slide, bullets: List[str], left: float, top: float,
                         width: float, height: float,
                         font_size: int = 15, color: Tuple = _C_DARK,
                         icon: str = "●") -> None:
    """Add a multi-line bullet list."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = _pt(4)
        run = p.add_run()
        short = textwrap.shorten(bullet, 110, placeholder="…")
        run.text = f"{icon}  {short}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = _rgb(color)


def _set_notes(slide, text: str) -> None:
    """Set presenter speaker notes on a slide."""
    if not text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _solid_bg(slide, color: Tuple) -> None:
    """Set a solid background color for the entire slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


# ── Individual slide builders ─────────────────────────────────────────────────

def _build_title_slide(prs, chapter_title: str, repo_name: str,
                        chapter_num: int, total_chapters: int, notes: str) -> None:
    """Slide 1 — Title slide with dark navy background."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _solid_bg(slide, _C_NAVY)

    # Top accent bar
    _add_rect(slide, 0, 0, _W_IN, 0.07, _C_BLUE)

    # Chapter number badge (large, ghost)
    badge = f"{chapter_num:02d}"
    _add_textbox(slide, badge, 0.4, 0.6, 3, 3,
                 font_size=130, bold=True, color=(255, 255, 255),
                 align_center=False)
    # Make it semi-transparent via alpha — workaround: use muted color
    txb = slide.shapes[-1]
    txb.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb((40, 70, 120))

    # Vertical accent line
    _add_rect(slide, 3.8, 1.0, 0.05, 5.0, _C_BLUE)

    # Chapter title
    _add_textbox(slide, chapter_title, 4.1, 1.3, 8.8, 1.8,
                 font_size=38, bold=True, color=_C_WHITE)

    # Repo name
    _add_textbox(slide, f"📦  {repo_name}", 4.1, 3.3, 8.8, 0.6,
                 font_size=18, color=(147, 197, 253))

    # Subtitle label
    _add_textbox(slide, "Code Tutorial", 4.1, 4.0, 4, 0.5,
                 font_size=14, color=(100, 150, 220))

    # Chapter position
    _add_textbox(slide, f"Chapter {chapter_num} of {total_chapters}",
                 4.1, 4.6, 4, 0.4, font_size=12, color=(80, 110, 160))

    # Date
    date_str = datetime.now().strftime("%B %d, %Y")
    _add_textbox(slide, date_str, 4.1, 5.1, 4, 0.4,
                 font_size=11, color=(80, 110, 160))

    # Bottom bar
    _add_rect(slide, 0, 7.3, _W_IN, 0.2, _C_BLUE)

    _set_notes(slide, notes)


def _build_overview_slide(prs, topic: str, bullets: List[str], notes: str) -> None:
    """Slide 2 — Overview / What is this?"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _solid_bg(slide, _C_OFFWHITE)

    # Left colored sidebar
    _add_rect(slide, 0, 0, 0.12, _H_IN, _C_BLUE)

    # Slide number badge
    _add_rect(slide, 0.3, 0.2, 0.7, 0.7, _C_BLUE)
    _add_textbox(slide, "2", 0.3, 0.2, 0.7, 0.7,
                 font_size=20, bold=True, color=_C_WHITE, align_center=True)

    # Heading
    heading = f"What is {topic}?"
    _add_textbox(slide, heading, 1.2, 0.15, 11.5, 0.9,
                 font_size=30, bold=True, color=_C_NAVY)

    # Underline
    _add_rect(slide, 1.2, 1.05, 6, 0.04, _C_BLUE)

    # Content card
    _add_rect(slide, 1.2, 1.2, 11.6, 5.8, _C_WHITE)

    # Bullets
    display_bullets = bullets[:5] if bullets else [
        f"{topic} is a core concept in software development.",
        "Understanding it helps you write better, more maintainable code.",
        "It is used widely across many projects and frameworks.",
    ]

    _add_textbox(slide, "Key Points:", 1.5, 1.35, 10.5, 0.4,
                 font_size=13, bold=True, color=_C_BLUE)

    _add_bullet_textbox(slide, display_bullets, 1.5, 1.85, 10.5, 4.8,
                         font_size=17, color=_C_DARK, icon="◆")

    # Bottom tag
    _add_rect(slide, 0, 7.25, _W_IN, 0.25, _C_NAVY)
    _add_textbox(slide, f"📚 CodeNarrator — {topic}", 0.3, 7.26, 10, 0.22,
                 font_size=11, color=(200, 220, 255))

    _set_notes(slide, notes)


def _build_code_slide(prs, topic: str, code: str, notes: str) -> None:
    """Slide 3 — Code Walkthrough with monospace code block."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _solid_bg(slide, _C_OFFWHITE)

    # Header bar
    _add_rect(slide, 0, 0, _W_IN, 1.0, _C_NAVY)
    _add_textbox(slide, "3", 0.3, 0.15, 0.7, 0.7,
                 font_size=20, bold=True, color=_C_BLUE, align_center=True)
    _add_textbox(slide, f"⌨  Code Walkthrough — {topic}",
                 1.2, 0.15, 11, 0.7, font_size=22, bold=True, color=_C_WHITE)

    if code:
        # Code block background
        _add_rect(slide, 0.3, 1.1, 8.2, 5.9, _C_CODE_BG)

        # Window bar dots
        _add_rect(slide, 0.5, 1.1, 7.8, 0.35, (30, 40, 60))
        for xi, dot_color in enumerate([(255, 95, 86), (255, 189, 46), (39, 201, 63)]):
            _add_rect(slide, 0.65 + xi * 0.3, 1.16, 0.18, 0.18, dot_color)
        _add_textbox(slide, f"  {topic.lower().replace(' ', '_')}.py",
                     1.4, 1.15, 6, 0.25, font_size=10, color=(150, 170, 200))

        # Code text
        code_lines = code[:1200]
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        code_box = slide.shapes.add_textbox(
            Inches(0.45), Inches(1.55), Inches(8.0), Inches(5.2)
        )
        tf = code_box.text_frame
        tf.word_wrap = False
        for i, line in enumerate(code_lines.splitlines()[:22]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line if line else " "
            run.font.name = "Consolas"
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*_C_CODE_TXT)

        # Explanation panel (right side)
        _add_rect(slide, 8.7, 1.1, 4.4, 5.9, (240, 245, 255))
        _add_textbox(slide, "📝 What this code does:", 8.85, 1.2, 4.1, 0.4,
                     font_size=12, bold=True, color=_C_NAVY)
        _add_rect(slide, 8.7, 1.65, 4.4, 0.03, _C_BLUE)

        # Extract key explanation lines from notes
        explanation = notes[:500] if notes else f"This code demonstrates key concepts of {topic}."
        _add_textbox(slide, explanation, 8.85, 1.75, 4.1, 5.0,
                     font_size=12, color=_C_DARK, word_wrap=True)
    else:
        # No code — show placeholder
        _add_rect(slide, 0.3, 1.1, 12.6, 5.9, (240, 245, 255))
        _add_textbox(slide, f"No code snippet available for this chapter.\n\nRefer to the chapter text for implementation details.",
                     0.5, 2.5, 12, 2, font_size=18, color=_C_MUTED, align_center=True)

    # Bottom bar
    _add_rect(slide, 0, 7.25, _W_IN, 0.25, _C_NAVY)
    _add_textbox(slide, f"📚 CodeNarrator — {topic}", 0.3, 7.26, 10, 0.22,
                 font_size=11, color=(200, 220, 255))

    _set_notes(slide, notes)


def _build_concepts_slide(prs, topic: str,
                           concepts: List[Tuple[str, str]], notes: str) -> None:
    """Slide 4 — Key Concepts grid."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _solid_bg(slide, _C_OFFWHITE)

    # Left sidebar
    _add_rect(slide, 0, 0, 0.12, _H_IN, _C_ACCENT)

    # Badge + heading
    _add_rect(slide, 0.3, 0.2, 0.7, 0.7, _C_ACCENT)
    _add_textbox(slide, "4", 0.3, 0.2, 0.7, 0.7,
                 font_size=20, bold=True, color=_C_WHITE, align_center=True)
    _add_textbox(slide, f"🔑  Key Concepts of {topic}",
                 1.2, 0.15, 11.5, 0.9, font_size=28, bold=True, color=_C_NAVY)
    _add_rect(slide, 1.2, 1.05, 6, 0.04, _C_ACCENT)

    # Fill up to 4 concepts
    display = concepts[:4]
    while len(display) < 4:
        display.append((f"Concept {len(display)+1}",
                        "A key idea related to this topic."))

    # 2x2 grid layout
    positions = [
        (0.3,  1.15),   # top-left
        (6.85, 1.15),   # top-right
        (0.3,  4.1),    # bottom-left
        (6.85, 4.1),    # bottom-right
    ]
    icons = ["🔹", "🔸", "🔷", "🔶"]
    accent_colors = [_C_BLUE, (245, 158, 11), (16, 185, 129), _C_ACCENT]

    for i, ((c_name, c_def), (cx, cy)) in enumerate(zip(display, positions)):
        # Card background
        _add_rect(slide, cx, cy, 6.3, 2.75, _C_WHITE)
        # Colored top strip
        _add_rect(slide, cx, cy, 6.3, 0.08, accent_colors[i])
        # Icon + name
        _add_textbox(slide, f"{icons[i]}  {c_name}", cx + 0.15, cy + 0.15,
                     5.9, 0.5, font_size=16, bold=True, color=accent_colors[i])
        # Definition
        short_def = textwrap.shorten(c_def, 140, placeholder="…")
        _add_textbox(slide, short_def, cx + 0.15, cy + 0.7, 5.9, 1.8,
                     font_size=13, color=_C_DARK, word_wrap=True)

    # Bottom bar
    _add_rect(slide, 0, 7.25, _W_IN, 0.25, _C_NAVY)
    _add_textbox(slide, f"📚 CodeNarrator — {topic}", 0.3, 7.26, 10, 0.22,
                 font_size=11, color=(200, 220, 255))

    _set_notes(slide, notes)


def _build_steps_slide(prs, topic: str, steps: List[str], notes: str) -> None:
    """Slide 5 — Tutorial Steps (numbered how-to list)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _solid_bg(slide, _C_OFFWHITE)

    # Header
    _add_rect(slide, 0, 0, _W_IN, 1.0, (16, 185, 129))  # green header
    _add_textbox(slide, "5", 0.3, 0.15, 0.7, 0.7,
                 font_size=20, bold=True, color=_C_WHITE, align_center=True)
    _add_textbox(slide, f"🛠  How to Use {topic} in Your Project",
                 1.2, 0.15, 11, 0.7, font_size=22, bold=True, color=_C_WHITE)

    # Steps (left column)
    _add_rect(slide, 0.3, 1.1, 7.8, 6.0, _C_WHITE)
    _add_textbox(slide, "Step-by-Step Guide:", 0.5, 1.18, 7.4, 0.4,
                 font_size=13, bold=True, color=_C_GREEN)

    display_steps = steps[:5] if steps else [
        f"Import and initialize {topic} in your project",
        "Configure the necessary parameters and settings",
        "Test the basic functionality with a simple example",
        "Integrate with the rest of your application",
        "Debug and optimize as needed",
    ]

    for i, step in enumerate(display_steps):
        step_y = 1.65 + i * 0.95
        # Number circle
        _add_rect(slide, 0.45, step_y, 0.5, 0.5, _C_GREEN)
        _add_textbox(slide, str(i + 1), 0.45, step_y, 0.5, 0.5,
                     font_size=14, bold=True, color=_C_WHITE, align_center=True)
        # Step text
        short = textwrap.shorten(step, 100, placeholder="…")
        _add_textbox(slide, short, 1.1, step_y + 0.03, 6.8, 0.55,
                     font_size=14, color=_C_DARK)

    # Tips panel (right column)
    _add_rect(slide, 8.4, 1.1, 4.6, 6.0, (240, 253, 244))
    _add_textbox(slide, "💡 Tips & Best Practices", 8.55, 1.2, 4.2, 0.45,
                 font_size=13, bold=True, color=_C_GREEN)
    _add_rect(slide, 8.4, 1.68, 4.6, 0.03, _C_GREEN)

    tip_text = (
        notes[:450] if notes
        else f"Always test {topic} with small inputs first.\n\n"
             f"Read the documentation carefully.\n\n"
             f"Look at existing examples in the codebase."
    )
    _add_textbox(slide, tip_text, 8.55, 1.78, 4.2, 5.1,
                 font_size=12, color=_C_DARK, word_wrap=True)

    # Bottom bar
    _add_rect(slide, 0, 7.25, _W_IN, 0.25, _C_NAVY)
    _add_textbox(slide, f"📚 CodeNarrator — {topic}", 0.3, 7.26, 10, 0.22,
                 font_size=11, color=(200, 220, 255))

    _set_notes(slide, notes)


def _build_summary_slide(prs, topic: str, summary_bullets: List[str],
                          next_chapter: str, notes: str) -> None:
    """Slide 6 — Summary + What's Next."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _solid_bg(slide, _C_NAVY)

    # Glow accent
    _add_rect(slide, 0, 0, _W_IN, 0.07, _C_BLUE)

    # Badge + heading
    _add_rect(slide, 0.4, 0.2, 0.7, 0.7, _C_BLUE)
    _add_textbox(slide, "6", 0.4, 0.2, 0.7, 0.7,
                 font_size=20, bold=True, color=_C_WHITE, align_center=True)
    _add_textbox(slide, f"✅  Summary — {topic}",
                 1.3, 0.2, 11, 0.7, font_size=28, bold=True, color=_C_WHITE)
    _add_rect(slide, 1.3, 0.97, 5, 0.04, _C_BLUE)

    # What you learned card
    _add_rect(slide, 0.4, 1.1, 7.8, 5.5, (20, 40, 80))
    _add_textbox(slide, "📋  What You Learned:", 0.6, 1.2, 7.4, 0.4,
                 font_size=14, bold=True, color=(147, 197, 253))

    disp_bullets = summary_bullets[:4] if summary_bullets else [
        f"The core purpose and function of {topic}",
        "How to read and understand the code",
        "Key concepts and their real-world applications",
        "How to implement this in your own project",
    ]
    for i, item in enumerate(disp_bullets):
        by = 1.72 + i * 0.85
        # Checkbox
        _add_rect(slide, 0.6, by, 0.42, 0.42, _C_GREEN)
        _add_textbox(slide, "✓", 0.6, by, 0.42, 0.42,
                     font_size=14, bold=True, color=_C_WHITE, align_center=True)
        short = textwrap.shorten(item, 90, placeholder="…")
        _add_textbox(slide, short, 1.15, by + 0.04, 6.8, 0.55,
                     font_size=15, color=_C_WHITE)

    # What's Next card
    _add_rect(slide, 8.5, 1.1, 4.5, 3.0, (20, 40, 80))
    _add_textbox(slide, "➡  What's Next:", 8.65, 1.2, 4.1, 0.4,
                 font_size=14, bold=True, color=(147, 197, 253))
    if next_chapter:
        _add_textbox(slide, next_chapter, 8.65, 1.7, 4.1, 1.8,
                     font_size=14, color=_C_WHITE, word_wrap=True)
    else:
        _add_textbox(slide, "Continue to the next chapter to build on these concepts.",
                     8.65, 1.7, 4.1, 1.8, font_size=14, color=_C_WHITE)

    # Practice card
    _add_rect(slide, 8.5, 4.3, 4.5, 2.3, (20, 40, 80))
    _add_textbox(slide, "🎯  Practice:", 8.65, 4.4, 4.1, 0.4,
                 font_size=14, bold=True, color=(147, 197, 253))
    _add_textbox(slide, f"Try implementing {topic} in a small personal project to reinforce your understanding.",
                 8.65, 4.9, 4.1, 1.5, font_size=13, color=_C_WHITE, word_wrap=True)

    # Bottom bar
    _add_rect(slide, 0, 7.25, _W_IN, 0.25, _C_BLUE)
    _add_textbox(slide, "🎓 Generated by CodeNarrator — Keep learning!", 0.3, 7.26, 12, 0.22,
                 font_size=11, color=_C_WHITE)

    _set_notes(slide, notes)


# ── Global cover slide ────────────────────────────────────────────────────────

def _build_cover_slide(prs, repo_name: str, chapter_count: int) -> None:
    """Opening cover slide for the whole tutorial."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _solid_bg(slide, _C_NAVY)

    _add_rect(slide, 0, 0, _W_IN, 0.1, _C_BLUE)
    _add_rect(slide, 0, 7.4, _W_IN, 0.1, _C_BLUE)

    # Logo area
    _add_rect(slide, 0.5, 0.8, 1.5, 1.5, _C_BLUE)
    _add_textbox(slide, "📚", 0.5, 0.8, 1.5, 1.5,
                 font_size=52, align_center=True, color=_C_WHITE)

    # Title
    _add_textbox(slide, repo_name, 2.4, 0.9, 10.5, 1.2,
                 font_size=42, bold=True, color=_C_WHITE)

    # Subtitle
    _add_textbox(slide, "AI-Generated Code Tutorial", 2.4, 2.2, 10, 0.6,
                 font_size=22, color=(147, 197, 253), italic=True)

    # Divider
    _add_rect(slide, 2.4, 2.95, 8, 0.05, _C_BLUE)

    # Stats
    date_str = datetime.now().strftime("%B %d, %Y")
    _add_textbox(slide, f"📖  {chapter_count} Chapters  ·  🎯  Beginner Friendly  ·  📅  {date_str}",
                 2.4, 3.2, 10, 0.5, font_size=15, color=(147, 197, 253))

    # Instructions
    _add_textbox(slide,
                 "This presentation was generated by CodeNarrator.\n"
                 "Each chapter contains 6 slides: Title · Overview · Code · Concepts · Steps · Summary\n"
                 "Speaker notes on every slide provide detailed explanations for self-study.",
                 2.4, 4.0, 10, 1.5, font_size=14, color=(180, 200, 240), word_wrap=True)

    # Badge
    _add_rect(slide, 2.4, 5.8, 3.5, 0.7, (40, 80, 150))
    _add_textbox(slide, "⚡  Generated by CodeNarrator", 2.4, 5.8, 3.5, 0.7,
                 font_size=14, bold=True, color=_C_WHITE, align_center=True)


# ── Main export function ───────────────────────────────────────────────────────

async def export_to_pptx_educational(
    job_output_dir: str,
) -> Dict[str, Any]:
    """
    Main async entry point. Reads markdown chapters from disk,
    builds a 6-slide-per-chapter educational PPTX with AI speaker notes.

    Returns:
        {"success": True,  "output_path": "...", "slides": N, "chapters": N}
        {"success": False, "error": "..."}
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return {"success": False,
                "error": "python-pptx not installed. Run: pip install python-pptx"}

    out_dir  = Path(job_output_dir)
    out_pptx = out_dir / "tutorial_educational.pptx"

    # ── Load markdown chapters ────────────────────────────────────────────────
    md_files = sorted(
        [f for f in out_dir.glob("*.md") if f.name != "index.md"],
        key=lambda p: p.name,
    )
    if not md_files:
        return {"success": False,
                "error": "No chapter .md files found — run text generation first"}

    # Derive repo name from index.md or folder
    repo_name = "Code Tutorial"
    index_path = out_dir / "index.md"
    if index_path.exists():
        idx_text = index_path.read_text(encoding="utf-8")
        m = re.search(r"^#\s+(.+)", idx_text, re.MULTILINE)
        if m:
            repo_name = m.group(1).strip()[:60]

    # ── Parse all chapters ────────────────────────────────────────────────────
    chapters = []
    for md_file in md_files:
        try:
            text = md_file.read_text(encoding="utf-8")
            parsed = _parse_chapter(text)
            chapters.append(parsed)
        except Exception as exc:
            logger.warning("Could not parse %s: %s", md_file.name, exc)

    if not chapters:
        return {"success": False, "error": "Could not parse any chapter files"}

    total_chapters = len(chapters)
    logger.info("PPTX educational: %d chapters", total_chapters)
    print(f"PPTXedu: generating {total_chapters} chapters × 6 slides + 1 cover")

    # ── Generate all speaker notes concurrently ───────────────────────────────
    # 5 notes per chapter × N chapters — run in parallel batches of 5
    print("PPTXedu: generating AI speaker notes…")

    async def _notes_for_chapter(ch: Dict) -> Dict[str, str]:
        topic = ch["title"]
        ctx   = ch["motivation"] or ch["core"] or ch["all_text"]
        results = await asyncio.gather(
            _gen_notes(topic, ctx,                 "overview"),
            _gen_notes(topic, ch["best_code"],     "code"),
            _gen_notes(topic, ch["core"],          "concepts"),
            _gen_notes(topic, ch["practical"],     "steps"),
            _gen_notes(topic, ch["conclusion"],    "summary"),
        )
        return {
            "title":    f"This chapter introduces {topic}. {results[0][:200]}",
            "overview": results[0],
            "code":     results[1],
            "concepts": results[2],
            "steps":    results[3],
            "summary":  results[4],
        }

    # Run chapter notes in batches to respect rate limits
    all_notes = []
    for i in range(0, total_chapters, 2):
        batch = chapters[i:i+2]
        batch_notes = await asyncio.gather(*[_notes_for_chapter(ch) for ch in batch])
        all_notes.extend(batch_notes)
        if i + 2 < total_chapters:
            await asyncio.sleep(2)   # brief pause between batches

    # ── Build presentation ────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(_W_IN)
    prs.slide_height = Inches(_H_IN)

    # Cover slide
    _build_cover_slide(prs, repo_name, total_chapters)

    slides_added = 1

    for i, (ch, notes) in enumerate(zip(chapters, all_notes), start=1):
        topic   = ch["title"]
        print(f"PPTXedu: [{i}/{total_chapters}] building 6 slides for '{topic}'")

        # Next chapter name for summary slide
        next_ch = chapters[i]["title"] if i < total_chapters else ""

        _build_title_slide(prs, topic, repo_name, i, total_chapters, notes["title"])
        _build_overview_slide(prs, topic, ch["overview_bullets"], notes["overview"])
        _build_code_slide(prs, topic, ch["best_code"], notes["code"])
        _build_concepts_slide(prs, topic, ch["key_concepts"], notes["concepts"])
        _build_steps_slide(prs, topic, ch["steps"], notes["steps"])
        _build_summary_slide(prs, topic, ch["summary_bullets"], next_ch, notes["summary"])

        slides_added += 6

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(str(out_pptx))
    size_mb = out_pptx.stat().st_size / (1024 * 1024)
    print(f"PPTXedu: DONE → tutorial_educational.pptx  {size_mb:.1f} MB  "
          f"{slides_added} slides  {total_chapters} chapters")

    return {
        "success":    True,
        "output_path": str(out_pptx),
        "slides":     slides_added,
        "chapters":   total_chapters,
        "size_mb":    round(size_mb, 2),
    }


# ── Backwards-compat wrapper (keeps old endpoint working) ─────────────────────

def export_to_pptx(job_output_dir: str, video_script: list) -> Dict[str, Any]:
    """
    Synchronous wrapper kept for backwards compatibility.
    Runs the async educational export via asyncio.run().
    video_script param is accepted but not used — chapters come from disk.
    """
    try:
        return asyncio.run(export_to_pptx_educational(job_output_dir))
    except RuntimeError:
        # Already inside an event loop
        loop = asyncio.get_event_loop()
        return loop.run_until_complete(export_to_pptx_educational(job_output_dir))
